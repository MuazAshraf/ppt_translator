import os
import shutil
import subprocess
import time
import zipfile
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional

import requests
from dotenv import load_dotenv
from openai import OpenAI
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt
from reportlab.pdfbase.pdfmetrics import stringWidth

load_dotenv('.env')
OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')
DEEPL_API_KEY = os.getenv('DEEPL_API_KEY')

if os.name == "nt":
    DEFAULT_SOFFICE = "soffice.com"
else:
    DEFAULT_SOFFICE = "soffice"

LIBREOFFICE_PATH = os.getenv("LIBREOFFICE_PATH") or DEFAULT_SOFFICE

ALLOWED_FORMATS = {"pptx", "pdf"}

LIBREOFFICE_CONVERT_OPTIONS = {
    "pdf": {
        "filters": ("pdf:impress_pdf_Export", "pdf"),
        "extensions": (".pdf",),
    }
}

LIBREOFFICE_FLAGS = [
    "--headless",
    "--invisible",
    "--nologo",
    "--nodefault",
    "--nolockcheck",
    "--norestore",
    "--nofirststartwizard",
    "--nocrashreport",
]

_libreoffice_checked = False
_libreoffice_available = False

client = OpenAI(api_key=OPENAI_API_KEY)

# Simple text measurement using ReportLab
def get_text_width(text, font_size, font_name="Helvetica"):
    """Get text width in points using ReportLab"""
    try:
        return stringWidth(text, font_name, font_size)
    except:
        # Fallback estimation
        return len(text) * font_size * 0.6

def calculate_font_size(text, max_width, current_size, min_size=9.0):
    """Calculate readable font size that fits - CONSERVATIVE approach"""
    if not text.strip():
        return current_size
    
    # Measure current text width
    current_width = get_text_width(text, current_size)
    
    # If it fits well, keep current size
    if current_width <= max_width * 1.1:  # Allow 10% overflow
        return current_size
    
    # CONSERVATIVE SCALING: Maximum 2-3 point reduction per step
    if current_size >= 20:
        max_reduction = 2.0  # Only reduce by 2pt max for large fonts
        conservative_min = current_size - max_reduction
    elif current_size >= 14:
        max_reduction = 1.5  # Reduce by 1.5pt max for medium fonts  
        conservative_min = current_size - max_reduction
    else:
        max_reduction = 1.0  # Reduce by 1pt max for small fonts
        conservative_min = current_size - max_reduction
    
    # Set reasonable minimums based on original size
    if current_size >= 24:
        absolute_min = 20.0  # Never below 20pt for titles
    elif current_size >= 16:
        absolute_min = 13.0  # Never below 13pt for body text
    else:
        absolute_min = max(min_size, 10.0)  # Never below 10pt for small text
    
    # Use the higher of conservative reduction or absolute minimum
    final_min = max(conservative_min, absolute_min)
    
    # Calculate scale factor but keep conservative
    scale_factor = max_width / current_width
    new_size = current_size * scale_factor
    
    # Apply conservative limits
    return max(new_size, final_min)

def translate_with_google(text, target_lang="es", source_lang="auto"):
    """Translate text using Google Translate API (free tier)"""
    if not text.strip():
        return text

    try:
        url = "https://translate.googleapis.com/translate_a/single"
        params = {
            'client': 'gtx',
            'sl': source_lang,
            'tl': target_lang,
            'dt': 't',
            'q': text
        }
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
        }
        response = requests.get(url, params=params, headers=headers, timeout=10)
        response.raise_for_status()
        result = response.json()
        if result and len(result) > 0 and result[0]:
            translated_text = ''.join([item[0] for item in result[0] if item[0]])
            return translated_text
        return text
    except Exception as e:
        print(f"Google Translate error: {e}")
        return text

def translate_with_deepl(text, target_lang="es", api_key=None):
    """Translate text using DeepL API"""
    if not text.strip() or not api_key:
        return text
    try:
        url = "https://api-free.deepl.com/v2/translate"
        data = {
            'auth_key': api_key,
            'text': text,
            'target_lang': target_lang.upper()
        }
        response = requests.post(url, data=data, timeout=10)
        response.raise_for_status()
        result = response.json()
        return result['translations'][0]['text']
    except Exception as e:
        print(f"DeepL error: {e}")
        return text

def translate_with_openai(text, target_lang="es", api_key=None):
    """Translate text using OpenAI GPT-4"""
    if not text.strip() or not api_key:
        return text
    try:
        client = OpenAI(api_key=api_key)
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": f"Translate to {target_lang}. Return ONLY the translation, no explanations."},
                {"role": "user", "content": text}
            ],
            temperature=0.3
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        print(f"OpenAI error: {e}")
        return text

def translate_text(text, target_lang="es", service="google", api_key=None):
    """Main translation function"""
    if service.lower() == "deepl":
        return translate_with_deepl(text, target_lang, api_key)
    elif service.lower() == "openai":
        return translate_with_openai(text, target_lang, api_key)
    elif service.lower() == "google":
        return translate_with_google(text, target_lang)
    else:
        print(f"Unknown service: {service}. Using Google Translate as fallback.")
    return translate_with_google(text, target_lang)

def translate_pptx(input_file, output_file, target_lang="es", service="google", api_key=None):
    """Simple PowerPoint translator with smart font sizing"""
    if not os.path.exists(input_file):
        print(f"Error: Input file '{input_file}' not found!")
        return 0

    print(f"Loading {input_file}...")
    try:
        prs = Presentation(input_file)
    except Exception as e:
        print(f"Error loading PowerPoint file: {e}")
        return 0

    total_translated = 0

    # Auto-load API keys
    if service.lower() == "deepl" and not api_key:
        api_key = DEEPL_API_KEY
    elif service.lower() == "openai" and not api_key:
        api_key = OPENAI_API_KEY

    print(f"Using {service.upper()} translation service...")

    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"Translating slide {slide_num}...")

        for shape in slide.shapes:
            # Handle regular text boxes
            if hasattr(shape, "text_frame") and shape.text_frame:
                # First pass: Translate all text
                original_size = 12.0
                all_translated_text = ""
                
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text and run.text.strip():
                            try:
                                # Capture original font size (use first run's size)
                                if hasattr(run.font, 'size') and run.font.size and original_size == 12.0:
                                    original_size = run.font.size.pt
                                
                                # Translate
                                original_text = run.text
                                translated_text = translate_text(original_text, target_lang, service, api_key)
                                run.text = translated_text
                                total_translated += 1
                                
                                # Collect all translated text for consistent sizing
                                all_translated_text += translated_text + " "
                                
                                time.sleep(0.1)
                            except Exception as e:
                                print(f"  Warning: Couldn't translate run text - {e}")
                
                # Second pass: Apply consistent font size to ALL runs
                if all_translated_text.strip() and shape.width:
                    shape_width = shape.width.inches * 72  # Convert to points
                    available_width = shape_width - 20  # Account for margins
                    
                    optimal_size = calculate_font_size(
                        all_translated_text.strip(), available_width, original_size, min_size=11.0  # Higher minimum for readability
                    )
                    
                    if abs(optimal_size - original_size) > 0.5:
                        # Apply same size to ALL runs in this text frame
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text.strip():
                                    run.font.size = Pt(optimal_size)
                        print(f"    Conservative font: {original_size:.1f}pt → {optimal_size:.1f}pt (entire text frame)")

                # Enable auto-fit
                try:
                    shape.text_frame.word_wrap = True
                    shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                except:
                    pass

            # Handle tables
            if hasattr(shape, 'has_table') and shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        if cell.text_frame:
                            # First pass: Translate all text in this cell
                            original_size = 10.0
                            all_cell_text = ""
                            
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.text and run.text.strip():
                                        try:
                                            # Capture original font size (use first run's size)
                                            if hasattr(run.font, 'size') and run.font.size and original_size == 10.0:
                                                original_size = run.font.size.pt
                                            
                                            # Translate
                                            original_text = run.text
                                            translated_text = translate_text(original_text, target_lang, service, api_key)
                                            run.text = translated_text
                                            total_translated += 1
                                            
                                            # Collect all translated text for consistent sizing
                                            all_cell_text += translated_text + " "
                                            
                                            time.sleep(0.1)
                                        except Exception as e:
                                            print(f"  Warning: Table cell translation failed - {e}")
                            
                            # Second pass: Apply consistent font size to ALL runs in this cell
                            if all_cell_text.strip():
                                optimal_size = calculate_font_size(
                                    all_cell_text.strip(), 100, original_size, min_size=10.0  # Higher minimum for table cells
                                )
                                
                                if abs(optimal_size - original_size) > 0.5:
                                    # Apply same size to ALL runs in this table cell
                                    for paragraph in cell.text_frame.paragraphs:
                                        for run in paragraph.runs:
                                            if run.text.strip():
                                                run.font.size = Pt(optimal_size)
                                    print(f"    Conservative table cell: {original_size:.1f}pt → {optimal_size:.1f}pt")

    print(f"Saving to {output_file}...")
    try:
        prs.save(output_file)
        print(f"Done! Translated {total_translated} elements using {service.upper()}")
        return total_translated
    except Exception as e:
        print(f"Error saving file: {e}")
        return 0

# LibreOffice helper functions
def _ensure_libreoffice_available():
    global _libreoffice_checked, _libreoffice_available
    if _libreoffice_checked:
        if not _libreoffice_available:
            raise FileNotFoundError("LibreOffice (soffice) not found.")
        return
    
    resolved = shutil.which(LIBREOFFICE_PATH)
    _libreoffice_checked = True
    _libreoffice_available = bool(resolved)
    if not _libreoffice_available:
        raise FileNotFoundError("LibreOffice (soffice) not found.")

def _normalize_formats(formats):
    if not formats:
        return ["pptx"]
    normalized = []
    for fmt in formats:
        if fmt:
            cleaned = fmt.lower().strip().lstrip(".")
            if cleaned and cleaned not in normalized:
                normalized.append(cleaned)
    if not normalized:
        normalized.append("pptx")
    invalid = set(normalized) - ALLOWED_FORMATS
    if invalid:
        raise ValueError(f"Unsupported formats: {sorted(invalid)}")
    return normalized

def _ensure_output_root(output_root, input_path):
    if output_root:
        root = Path(output_root)
    else:
        root = input_path.parent / f"{input_path.stem}_translations"
    root.mkdir(parents=True, exist_ok=True)
    return root

def _libreoffice_convert(source_path, target_format):
    options = LIBREOFFICE_CONVERT_OPTIONS.get(target_format)
    if not options:
        raise ValueError(f"Unsupported format: {target_format}")
    
    _ensure_libreoffice_available()
    filters = options["filters"]
    extensions = options["extensions"]
    outdir = source_path.parent
    
    for convert_arg in filters:
        cmd = [LIBREOFFICE_PATH, *LIBREOFFICE_FLAGS, "--convert-to", convert_arg, "--outdir", str(outdir), str(source_path)]
        result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
        if result.returncode == 0:
            break
    else:
        raise RuntimeError("LibreOffice conversion failed")
    
    # Find produced file
    for _ in range(10):
        for extension in extensions:
            pattern = f"{source_path.stem}*{extension}"
            candidates = list(source_path.parent.glob(pattern))
        if candidates:
                return max(candidates, key=lambda p: p.stat().st_mtime)
        time.sleep(0.2)

    raise RuntimeError(f"No {target_format} file was produced")

def _bundle_outputs_to_zip(output_root):
    zip_path = output_root.parent / f"{output_root.name}.zip"
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as archive:
        for file_path in output_root.rglob("*"):
            if file_path.is_file():
                archive.write(file_path, arcname=file_path.relative_to(output_root))
    return zip_path

def translate_pptx_multi(input_file, target_langs, service="google", api_key=None, 
                        formats=None, output_root=None, max_workers=None, zip_output=True):
    """Multi-language PowerPoint translator"""
    input_path = Path(input_file)
    if not input_path.exists():
        raise FileNotFoundError(f"Input file '{input_path}' not found.")
    
    languages = [lang.strip() for lang in target_langs if lang and lang.strip()]
    if not languages:
        raise ValueError("At least one valid target language must be provided.")
    
    normalized_formats = _normalize_formats(formats)
    if "pdf" in normalized_formats:
            _ensure_libreoffice_available()
    
    if service.lower() == "deepl" and not api_key:
        api_key = DEEPL_API_KEY
    elif service.lower() == "openai" and not api_key:
        api_key = OPENAI_API_KEY
    
    output_root_path = _ensure_output_root(output_root, input_path)
    
    print(f"🚀 Starting translation to {len(languages)} languages...")
    
    translations = {}
    errors = []
    
    for language in languages:
        try:
            print(f"\n📝 Translating to {language.upper()}...")
            lang_dir = output_root_path / language
            lang_dir.mkdir(exist_ok=True)
            pptx_path = lang_dir / f"{input_path.stem}_{language}.pptx"
            
            # Translate
            translated = translate_pptx(str(input_path), str(pptx_path), language, service, api_key)
            
            outputs = {}
            if "pptx" in normalized_formats:
                outputs["pptx"] = str(pptx_path)
            
            if "pdf" in normalized_formats and pptx_path.exists():
                print(f"  📄 Converting to PDF...")
                try:
                    pdf_path = _libreoffice_convert(pptx_path, "pdf")
                    outputs["pdf"] = str(pdf_path)
                    print(f"  ✅ PDF created: {pdf_path}")
                except Exception as exc:
                    errors.append(f"{language} PDF export failed: {exc}")
            
            translations[language] = {"count": translated, "outputs": outputs}
            print(f"✅ {language.upper()}: {translated} elements translated")
            
        except Exception as exc:
            errors.append(f"{language}: {exc}")
            print(f"❌ {language.upper()}: {exc}")
    
    zip_path = None
    if zip_output and translations:
        try:
            print("\n📦 Creating ZIP archive...")
            zip_path = _bundle_outputs_to_zip(output_root_path)
            print(f"✅ ZIP created: {zip_path}")
        except Exception as exc:
            errors.append(f"ZIP packaging failed: {exc}")
    
    return {
        "input": str(input_path),
        "output_root": str(output_root_path),
        "zip_path": str(zip_path) if zip_path else None,
        "formats": normalized_formats,
        "translations": translations,
        "errors": errors,
    }